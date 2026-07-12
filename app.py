"""ICU餐厅会员看板 — Flask Application."""
import sqlite3
import os
from functools import wraps
from flask import Flask, request, jsonify, render_template, redirect, url_for, session
import bcrypt

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", os.urandom(24).hex())

DB_PATH = os.path.join(os.path.dirname(__file__), "data.db")


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA busy_timeout=5000")  # 5秒超时，多用户并发不锁死
    return conn


def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if "user_id" not in session:
            if request.is_json or request.path.startswith("/api/"):
                return jsonify({"error": "未登录"}), 401
            return redirect(url_for("login_page"))
        return f(*args, **kwargs)
    return decorated


# ═══════════════════════════════════════════
# Page Routes
# ═══════════════════════════════════════════

@app.route("/")
def index():
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    return render_template("dashboard.html")


@app.route("/login")
def login_page():
    return render_template("login.html")


# ═══════════════════════════════════════════
# Auth APIs
# ═══════════════════════════════════════════

@app.route("/api/login", methods=["POST"])
def api_login():
    data = request.get_json()
    username = data.get("username", "").strip()
    password = data.get("password", "")

    conn = get_db()
    user = conn.execute("SELECT * FROM users WHERE username = ?", (username,)).fetchone()
    conn.close()

    if not user or not bcrypt.checkpw(password.encode(), user["password_hash"].encode()):
        return jsonify({"error": "用户名或密码错误"}), 401

    session["user_id"] = user["id"]
    session["username"] = user["username"]
    session["display_name"] = user["display_name"] or user["username"]
    return jsonify({"ok": True, "display_name": session["display_name"]})


@app.route("/api/logout", methods=["POST"])
def api_logout():
    session.clear()
    return jsonify({"ok": True})


@app.route("/api/me")
@login_required
def api_me():
    return jsonify({
        "username": session.get("username"),
        "display_name": session.get("display_name"),
    })


# ═══════════════════════════════════════════
# Overview APIs
# ═══════════════════════════════════════════

@app.route("/api/overview")
@login_required
def api_overview():
    """Full overview dashboard."""
    conn = get_db()
    latest_month = conn.execute("SELECT MAX(月份) FROM restaurants").fetchone()[0]

    # 1. Summary stats
    summary = conn.execute("""
        SELECT
            COUNT(DISTINCT 门店名称) as total_stores,
            COUNT(DISTINCT 城市) as total_cities,
            COUNT(DISTINCT 战队) as total_teams
        FROM restaurants WHERE 月份 = ?
    """, (latest_month,)).fetchone()

    # 2. Team comparison (latest month)
    teams = conn.execute("""
        SELECT
            战队,
            COUNT(DISTINCT 门店名称) as store_count,
            ROUND(AVG(月消费人数), 0) as avg_consumers,
            ROUND(AVG(人均月频次), 2) as avg_frequency,
            ROUND(AVG(COALESCE(超会消费占比, 0)), 4) as avg_super_ratio,
            ROUND(AVG(COALESCE(好友消费占比, 0)), 4) as avg_friend_ratio,
            ROUND(AVG(COALESCE(超会售卡, 0)), 1) as avg_card_sales,
            AVG(COALESCE(堂食日均, 0)) as avg_dinein_daily,
            SUM(CASE WHEN 机会点_消费人数不足 = 1 THEN 1 ELSE 0 END) as opp_consumer,
            SUM(CASE WHEN 机会点_人均频次不足 = 1 THEN 1 ELSE 0 END) as opp_frequency,
            SUM(CASE WHEN 机会点_超会不足 = 1 THEN 1 ELSE 0 END) as opp_super,
            SUM(CASE WHEN 机会点_好友拉新不足 = 1 THEN 1 ELSE 0 END) as opp_friend
        FROM restaurants
        WHERE 月份 = ?
        GROUP BY 战队
        ORDER BY avg_consumers DESC
    """, (latest_month,)).fetchall()

    teams_data = []
    for t in teams:
        total = t["store_count"]
        teams_data.append({
            "name": t["战队"],
            "store_count": total,
            "avg_consumers": t["avg_consumers"],
            "avg_frequency": t["avg_frequency"],
            "avg_super_ratio": t["avg_super_ratio"],
            "avg_friend_ratio": t["avg_friend_ratio"],
            "avg_card_sales": t["avg_card_sales"],
            "avg_dinein_daily": round(t["avg_dinein_daily"], 0) if t["avg_dinein_daily"] else None,
            "opp_consumer_pct": round(t["opp_consumer"] / total * 100, 1) if total else 0,
            "opp_frequency_pct": round(t["opp_frequency"] / total * 100, 1) if total else 0,
            "opp_super_pct": round(t["opp_super"] / total * 100, 1) if total else 0,
            "opp_friend_pct": round(t["opp_friend"] / total * 100, 1) if total else 0,
        })

    # 3. City ranking (top 20 by consumers)
    cities = conn.execute("""
        SELECT
            城市,
            COUNT(DISTINCT 门店名称) as store_count,
            ROUND(AVG(月消费人数), 0) as avg_consumers,
            ROUND(AVG(人均月频次), 2) as avg_frequency,
            ROUND(AVG(COALESCE(堂食日均, 0)), 0) as avg_dinein_daily
        FROM restaurants
        WHERE 月份 = ?
        GROUP BY 城市
        ORDER BY avg_consumers DESC
        LIMIT 20
    """, (latest_month,)).fetchall()

    cities_data = [{
        "name": c["城市"], "store_count": c["store_count"],
        "avg_consumers": c["avg_consumers"], "avg_frequency": c["avg_frequency"],
        "avg_dinein_daily": c["avg_dinein_daily"],
    } for c in cities]

    # 4. Monthly trend (overall group)
    months_rows = conn.execute("""
        SELECT 月份,
            COUNT(DISTINCT 门店名称) as store_count,
            ROUND(AVG(月消费人数), 0) as avg_consumers,
            ROUND(AVG(人均月频次), 2) as avg_frequency,
            ROUND(AVG(COALESCE(超会消费占比, 0)), 4) as avg_super_ratio,
            ROUND(AVG(COALESCE(好友消费占比, 0)), 4) as avg_friend_ratio,
            ROUND(AVG(COALESCE(堂食日均, 0)), 0) as avg_dinein_daily
        FROM restaurants
        GROUP BY 月份
        ORDER BY 月份
    """).fetchall()

    trend = {
        "months": [m["月份"] for m in months_rows],
        "store_count": [m["store_count"] for m in months_rows],
        "avg_consumers": [m["avg_consumers"] for m in months_rows],
        "avg_frequency": [m["avg_frequency"] for m in months_rows],
        "avg_super_ratio": [m["avg_super_ratio"] for m in months_rows],
        "avg_friend_ratio": [m["avg_friend_ratio"] for m in months_rows],
        "avg_dinein_daily": [m["avg_dinein_daily"] for m in months_rows],
    }

    # 5. Customer structure: new vs old vs returning by month (from city_new_old)
    cust_struct = conn.execute("""
        SELECT
            月份,
            AVG(店均顾客数) as avg_customers,
            AVG(新客数) as avg_new,
            AVG(活跃客数) as avg_active,
            AVG(回流客数) as avg_returning,
            AVG(新客复购率) as avg_new_repurchase,
            AVG(老客回流率) as avg_old_return,
            COUNT(*) as city_count
        FROM city_new_old
        GROUP BY 月份
        ORDER BY 月份
    """).fetchall()

    # Lost customers per store = city_lost / store_count per city
    lost_struct = conn.execute("""
        SELECT cl.月份,
            ROUND(AVG(cl.流失客数 * 1.0 / rs.store_cnt), 0) as avg_lost
        FROM city_lost cl
        INNER JOIN (
            SELECT 月份, 城市, COUNT(DISTINCT 门店名称) as store_cnt
            FROM restaurants GROUP BY 月份, 城市
        ) rs ON cl.月份 = rs.月份 AND cl.城市 = rs.城市
        GROUP BY cl.月份
        ORDER BY cl.月份
    """).fetchall()
    lost_map = {r["月份"]: r["avg_lost"] for r in lost_struct}

    cust_data = {
        "months": [c["月份"] for c in cust_struct],
        "avg_customers": [round(c["avg_customers"], 0) if c["avg_customers"] else 0 for c in cust_struct],
        "new_customers": [round(c["avg_new"], 0) if c["avg_new"] else 0 for c in cust_struct],
        "active_customers": [round(c["avg_active"], 0) if c["avg_active"] else 0 for c in cust_struct],
        "returning_customers": [round(c["avg_returning"], 0) if c["avg_returning"] else 0 for c in cust_struct],
        "lost_customers": [round(lost_map.get(c["月份"], 0), 0) for c in cust_struct],
        "new_repurchase_rate": [c["avg_new_repurchase"] for c in cust_struct],
        "old_return_rate": [c["avg_old_return"] for c in cust_struct],
        "city_count": [c["city_count"] for c in cust_struct],
    }

    # 6. Team trend over months
    team_trend = conn.execute("""
        SELECT 月份, 战队, ROUND(AVG(月消费人数), 0) as avg_consumers
        FROM restaurants
        GROUP BY 月份, 战队
        ORDER BY 月份, avg_consumers DESC
    """).fetchall()

    from collections import defaultdict
    team_trend_data = defaultdict(dict)
    all_teams_set = set()
    for r in team_trend:
        team_trend_data[r["月份"]][r["战队"]] = r["avg_consumers"]
        all_teams_set.add(r["战队"])

    # Sort teams by latest month consumers
    latest_team_ranking = sorted(teams_data, key=lambda x: x["avg_consumers"] or 0, reverse=True)
    ranked_teams = [t["name"] for t in latest_team_ranking]

    team_trend_series = {}
    for team in ranked_teams:
        team_trend_series[team] = [team_trend_data.get(m, {}).get(team) for m in trend["months"]]

    conn.close()

    return jsonify({
        "latest_month": latest_month,
        "total_stores": summary["total_stores"],
        "total_cities": summary["total_cities"],
        "total_teams": summary["total_teams"],
        "teams": teams_data,
        "cities": cities_data,
        "trend": trend,
        "customer_structure": cust_data,
        "team_trend": {
            "months": trend["months"],
            "teams": ranked_teams,
            "series": team_trend_series,
        },
    })


# ═══════════════════════════════════════════
# Store APIs
# ═══════════════════════════════════════════

@app.route("/api/stores/search")
@login_required
def api_store_search():
    q = request.args.get("q", "").strip()
    if not q:
        return jsonify([])

    conn = get_db()
    rows = conn.execute("""
        SELECT 餐厅编码, 门店名称, 城市, 战队, 区域, 城市类型, 商圈类型
        FROM restaurants
        WHERE 门店名称 LIKE ? OR 城市 LIKE ? OR 餐厅编码 LIKE ?
        GROUP BY 门店名称, 城市
        ORDER BY 门店名称
        LIMIT 30
    """, (f"%{q}%", f"%{q}%", f"%{q}%")).fetchall()
    conn.close()

    return jsonify([{
        "code": r["餐厅编码"] or "",
        "name": r["门店名称"], "city": r["城市"],
        "team": r["战队"], "region": r["区域"],
        "city_type": r["城市类型"], "area_type": r["商圈类型"],
    } for r in rows])


@app.route("/api/stores/list")
@login_required
def api_store_list():
    team = request.args.get("team", "").strip()
    city = request.args.get("city", "").strip()

    conn = get_db()
    query = "SELECT DISTINCT 门店名称, 城市, 战队, 餐厅编码 FROM restaurants WHERE 1=1"
    params = []
    if team:
        query += " AND 战队 = ?"
        params.append(team)
    if city:
        query += " AND 城市 = ?"
        params.append(city)
    query += " ORDER BY 城市, 门店名称"
    rows = conn.execute(query, params).fetchall()
    conn.close()

    return jsonify([{
        "name": r["门店名称"], "city": r["城市"],
        "team": r["战队"], "code": r["餐厅编码"]
    } for r in rows])


@app.route("/api/filters")
@login_required
def api_filters():
    conn = get_db()
    teams = [r[0] for r in conn.execute(
        "SELECT DISTINCT 战队 FROM restaurants WHERE 战队 IS NOT NULL ORDER BY 战队").fetchall()]
    cities = [r[0] for r in conn.execute(
        "SELECT DISTINCT 城市 FROM restaurants WHERE 城市 IS NOT NULL ORDER BY 城市").fetchall()]
    latest_month = conn.execute("SELECT MAX(月份) FROM restaurants").fetchone()[0]
    all_months = [r[0] for r in conn.execute(
        "SELECT DISTINCT 月份 FROM restaurants ORDER BY 月份").fetchall()]
    conn.close()
    return jsonify({"teams": teams, "cities": cities, "latest_month": latest_month, "months": all_months})


@app.route("/api/store/<store_name>/detail")
@login_required
def api_store_detail(store_name):
    conn = get_db()
    row = conn.execute("""
        SELECT * FROM restaurants
        WHERE 门店名称 = ?
        ORDER BY 月份 DESC LIMIT 1
    """, (store_name,)).fetchone()

    if not row:
        conn.close()
        return jsonify({"error": "未找到该餐厅"}), 404

    result = dict(row)
    city, city_type, month = row["城市"], row["城市类型"], row["月份"]

    peer = conn.execute("""
        SELECT
            AVG(月消费人数) as avg_consumers,
            AVG(人均月频次) as avg_frequency,
            AVG(超会消费占比) as avg_super_ratio,
            AVG(好友消费占比) as avg_friend_ratio,
            AVG(超会售卡) as avg_card_sales,
            COUNT(*) as peer_count
        FROM restaurants
        WHERE 城市类型 = ? AND 月份 = ? AND 门店名称 != ?
    """, (city_type, month, store_name)).fetchone()

    if peer and peer["peer_count"]:
        result["peer_consumers"] = round(peer["avg_consumers"], 1) if peer["avg_consumers"] else None
        result["peer_frequency"] = round(peer["avg_frequency"], 2) if peer["avg_frequency"] else None
        result["peer_super_ratio"] = round(peer["avg_super_ratio"], 4) if peer["avg_super_ratio"] else None
        result["peer_friend_ratio"] = round(peer["avg_friend_ratio"], 4) if peer["avg_friend_ratio"] else None
        result["peer_card_sales"] = round(peer["avg_card_sales"], 1) if peer["avg_card_sales"] else None
        result["peer_count"] = peer["peer_count"]

    # City-level data for this store's city
    city_data = conn.execute("""
        SELECT * FROM city_new_old WHERE 城市 = ? ORDER BY 月份 DESC LIMIT 1
    """, (city,)).fetchone()
    if city_data:
        result["city_data"] = dict(city_data)

    conn.close()
    return jsonify(result)


@app.route("/api/store/<store_name>/trend")
@login_required
def api_store_trend(store_name):
    conn = get_db()
    rows = conn.execute("""
        SELECT 月份, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比,
               超会售卡, 好友新增, 消费人数_同类均值, 频次_同类均值,
               超会_同类均值, 好友_同类均值,
               机会点_消费人数不足, 机会点_人均频次不足,
               机会点_超会不足, 机会点_好友拉新不足, 堂食日均
        FROM restaurants
        WHERE 门店名称 = ?
        ORDER BY 月份
    """, (store_name,)).fetchall()
    conn.close()

    if not rows:
        return jsonify({"error": "无数据"}), 404

    trend = {
        "months": [], "consumers": [], "peer_consumers": [],
        "frequency": [], "peer_frequency": [],
        "super_ratio": [], "peer_super_ratio": [],
        "friend_ratio": [], "peer_friend_ratio": [],
        "card_sales": [], "friend_new": [], "dinein_daily": [],
        "opportunity": {"consumers": [], "frequency": [], "super": [], "friend": []}
    }

    for r in rows:
        trend["months"].append(r["月份"])
        trend["consumers"].append(r["月消费人数"])
        trend["peer_consumers"].append(r["消费人数_同类均值"])
        trend["frequency"].append(r["人均月频次"])
        trend["peer_frequency"].append(r["频次_同类均值"])
        trend["super_ratio"].append(r["超会消费占比"])
        trend["peer_super_ratio"].append(r["超会_同类均值"])
        trend["friend_ratio"].append(r["好友消费占比"])
        trend["peer_friend_ratio"].append(r["好友_同类均值"])
        trend["card_sales"].append(r["超会售卡"])
        trend["friend_new"].append(r["好友新增"])
        trend["dinein_daily"].append(r["堂食日均"])
        trend["opportunity"]["consumers"].append(r["机会点_消费人数不足"])
        trend["opportunity"]["frequency"].append(r["机会点_人均频次不足"])
        trend["opportunity"]["super"].append(r["机会点_超会不足"])
        trend["opportunity"]["friend"].append(r["机会点_好友拉新不足"])

    return jsonify(trend)


@app.route("/api/store/<store_name>/ranking")
@login_required
def api_store_ranking(store_name):
    conn = get_db()
    store = conn.execute("""
        SELECT 城市类型, 月份 FROM restaurants
        WHERE 门店名称 = ? ORDER BY 月份 DESC LIMIT 1
    """, (store_name,)).fetchone()

    if not store:
        conn.close()
        return jsonify({"error": "未找到"}), 404

    city_type, month = store["城市类型"], store["月份"]

    peers = conn.execute("""
        SELECT 门店名称, 城市, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比, 堂食日均
        FROM restaurants
        WHERE 城市类型 = ? AND 月份 = ?
    """, (city_type, month)).fetchall()
    conn.close()

    total = len(peers)

    def get_rank(field, desc=True):
        sorted_peers = sorted(peers, key=lambda x: x[field] or 0, reverse=desc)
        for i, p in enumerate(sorted_peers):
            if p["门店名称"] == store_name:
                return i + 1
        return None

    return jsonify({
        "total_peers": total,
        "city_type": city_type,
        "rank_consumers": get_rank("月消费人数"),
        "rank_frequency": get_rank("人均月频次"),
        "rank_super": get_rank("超会消费占比"),
        "rank_friend": get_rank("好友消费占比"),
        "rank_dinein": get_rank("堂食日均"),
    })


# ═══════════════════════════════════════════
# Summary APIs (rule-based analysis)
# ═══════════════════════════════════════════

def generate_store_summary(store_name):
    conn = get_db()
    rows = conn.execute("""
        SELECT 月份, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比, 好友新增,
               消费人数_同类均值, 频次_同类均值, 超会_同类均值, 好友_同类均值,
               机会点_消费人数不足, 机会点_人均频次不足, 机会点_超会不足, 机会点_好友拉新不足,
               超会售卡, 周四平均售卡
        FROM restaurants WHERE 门店名称 = ? ORDER BY 月份
    """, (store_name,)).fetchall()

    city = conn.execute("SELECT 城市, 城市类型, 区域, 战队 FROM restaurants WHERE 门店名称 = ? ORDER BY 月份 DESC LIMIT 1",
                       (store_name,)).fetchone()
    # 复购率从 city_new_old 表里取（同店按城市）
    new_repurchase = None
    old_repurchase = None
    if city:
        crow = conn.execute("""
            SELECT AVG(新客复购率) as nr, AVG(老客回流率) as or_
            FROM city_new_old WHERE 城市 = ? ORDER BY 月份 DESC LIMIT 1
        """, (city['城市'],)).fetchone()
        if crow:
            new_repurchase = crow['nr']
            old_repurchase = crow['or_']
    conn.close()

    if not rows:
        return None

    latest = rows[-1]
    prev = rows[-2] if len(rows) >= 2 else None

    # 4维度诊断
    val_c = latest["月消费人数"];  peer_c = latest["消费人数_同类均值"]
    val_f = latest["人均月频次"];  peer_f = latest["频次_同类均值"]
    val_s = latest["超会消费占比"];peer_s = latest["超会_同类均值"]
    val_h = latest["好友消费占比"];peer_h = latest["好友_同类均值"]
    new_repurchase = new_repurchase
    old_repurchase = old_repurchase
    friend_new = latest["好友新增"]
    card_sales = latest["超会售卡"]
    thursday_card = latest["周四平均售卡"]

    # 格式化助手
    def fmt_int(v):
        if v is None: return "—"
        return f"{int(round(v)):,}"
    def fmt_pct(v):
        if v is None: return "—"
        return f"{v*100:.1f}%"
    def fmt_f(v):
        if v is None: return "—"
        return f"{v:.2f}"

    # 数据表现（拼数字）
    def desc(name, val, peer, fmt):
        if val is None: return f"{name}数据缺失"
        diff = (val - peer) if peer else 0
        sign = "低于" if diff < 0 else ("高于" if diff > 0 else "持平")
        return f"{name}{fmt(val)}（同类{fmt(peer)}{sign}）"

    data_issues = []
    if latest["机会点_消费人数不足"]:
        data_issues.append(f"月消费人数{desc('月消费人数', val_c, peer_c, fmt_int)}")
    if latest["机会点_人均频次不足"]:
        data_issues.append(f"月频次{desc('人均月频次', val_f, peer_f, fmt_f)}次（同类{fmt_f(peer_f)}次）")
        if new_repurchase is not None:
            data_issues[-1] += f"【堂食顾客复购率{fmt_pct(old_repurchase)}、新客复购率{fmt_pct(new_repurchase)}】"
    if latest["机会点_超会不足"]:
        data_issues.append(f"超会消费占比{desc('超会消费占比', val_s, peer_s, fmt_pct)}，月售卡{fmt_int(card_sales)}张")
    if latest["机会点_好友拉新不足"]:
        data_issues.append(f"好友消费占比{desc('好友消费占比', val_h, peer_h, fmt_pct)}，月新增好友{fmt_int(friend_new)}人")

    # 数字编号 ①②③④ — 严格按业务顺序而非"出现的顺序"
    issue_label_map = {
        0: '消费人数',
        1: '消费频次',
        2: '付费会员',
        3: '企微',
    }
    # 按业务顺序重排（保持①②③④编号稳定）
    business_order = [
        ('消费人数', '月消费人数'),
        ('消费频次', '月频次'),
        ('付费会员', '超会消费占比'),
        ('企微',     '好友消费占比'),
    ]
    business_filled = []
    for cat, key in business_order:
        for item in data_issues:
            if item.startswith(key):
                business_filled.append(item)
                break
    # 数字编号
    def num_cn(i):
        return {0:'①',1:'②',2:'③',3:'④',4:'⑤',5:'⑥'}.get(i, str(i+1))
    if business_filled:
        data_issues_numbered = [f"{num_cn(i)}{x}" for i, x in enumerate(business_filled)]
        data_line = "1.数据诊断：" + "；".join(data_issues_numbered)
    else:
        data_line = "1.数据诊断：各项指标均达到或优于同类均值水平。"

    # 改善动作建议
    actions = []
    if latest["机会点_消费人数不足"]:
        actions.append("加强周边渗透&引流揽客动作")
    if latest["机会点_人均频次不足"]:
        actions.append("安排专人推售超级会员，提升复购；调研顾客体验提升复购率")
    if latest["机会点_超会不足"]:
        actions.append(f"加强超会转化（当前月售卡{fmt_int(card_sales)}张，需提升）")
    if latest["机会点_好友拉新不足"]:
        actions.append(f"加强企微拉新（当前好友{fmt_int(friend_new)}人）")

    # 阶段改善目标（仅在「超会售卡」和「企微好友拉新」上给建议）
    # 规则：建议不超过当前水平的 150%；若已超过同类均值，则建议不超过 120%
    targets = []
    if latest["机会点_超会不足"] and card_sales:
        already_above_peer = (peer_s is not None) and (val_s is not None) and (val_s >= peer_s)
        ratio = 1.2 if already_above_peer else 1.5
        targets.append(f"超会月售卡目标{card_sales*ratio:.2f}张（当前{fmt_int(card_sales)}张）")
    if latest["机会点_好友拉新不足"] and friend_new:
        already_above_peer_h = (peer_h is not None) and (val_h is not None) and (val_h >= peer_h)
        ratio_h = 1.2 if already_above_peer_h else 1.5
        targets.append(f"企微拉新目标{friend_new*ratio_h:.2f}人（当前{fmt_int(friend_new)}人）")

    sentences = []
    sentences.append(data_line)

    if actions:
        action_line = "2.改善动作建议：" + "；".join(actions) + "。"
    else:
        action_line = "2.改善动作建议：保持当前节奏，关注数据变化。"
    sentences.append(action_line)

    if targets:
        target_line = "3.一阶段改善目标建议：" + "；".join(targets) + "。"
    else:
        target_line = "3.一阶段改善目标建议：巩固现有优势，争取同行标杆水平。"
    sentences.append(target_line)

    return {
        "sentences": sentences,
        "opp_count": len(data_issues),
        "opp_dimensions": [
            "消费人数" if latest["机会点_消费人数不足"] else None,
            "消费频次" if latest["机会点_人均频次不足"] else None,
            "超会" if latest["机会点_超会不足"] else None,
            "好友拉新" if latest["机会点_好友拉新不足"] else None,
        ],
    }


@app.route("/api/store/<store_name>/summary")
@login_required
def api_store_summary(store_name):
    result = generate_store_summary(store_name)
    if not result:
        return jsonify({"error": "未找到该餐厅"}), 404
    return jsonify(result)


# ═══════════════════════════════════════════
# Team & City APIs
# ═══════════════════════════════════════════

@app.route("/api/team/<team_name>/overview")
@login_required
def api_team_overview(team_name):
    conn = get_db()
    latest_month = conn.execute("SELECT MAX(月份) FROM restaurants").fetchone()[0]

    stores = conn.execute("""
        SELECT 门店名称, 城市, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比,
               机会点_消费人数不足, 机会点_人均频次不足, 机会点_超会不足, 机会点_好友拉新不足
        FROM restaurants
        WHERE 战队 = ? AND 月份 = ?
        ORDER BY 月消费人数 DESC
    """, (team_name, latest_month)).fetchall()

    trend = conn.execute("""
        SELECT 月份,
            ROUND(AVG(月消费人数), 0) as c,
            ROUND(AVG(人均月频次), 2) as f,
            ROUND(AVG(COALESCE(超会消费占比, 0)), 4) as sr,
            ROUND(AVG(COALESCE(好友消费占比, 0)), 4) as fr
        FROM restaurants
        WHERE 战队 = ?
        GROUP BY 月份 ORDER BY 月份
    """, (team_name,)).fetchall()

    conn.close()

    total = len(stores)
    opp_counts = sum(1 for s in stores if s["机会点_消费人数不足"])
    opp_freq = sum(1 for s in stores if s["机会点_人均频次不足"])
    opp_super = sum(1 for s in stores if s["机会点_超会不足"])
    opp_friend = sum(1 for s in stores if s["机会点_好友拉新不足"])

    return jsonify({
        "team": team_name,
        "store_count": total,
        "stores": [dict(s) for s in stores],
        "trend": {
            "months": [t["月份"] for t in trend],
            "consumers": [t["c"] for t in trend],
            "frequency": [t["f"] for t in trend],
            "super_ratio": [t["sr"] for t in trend],
            "friend_ratio": [t["fr"] for t in trend],
        },
        "opportunity": {
            "consumer_pct": round(opp_counts / total * 100, 1) if total else 0,
            "frequency_pct": round(opp_freq / total * 100, 1) if total else 0,
            "super_pct": round(opp_super / total * 100, 1) if total else 0,
            "friend_pct": round(opp_friend / total * 100, 1) if total else 0,
        }
    })


@app.route("/api/city/<city_name>/overview")
@login_required
def api_city_overview(city_name):
    conn = get_db()
    latest_month = conn.execute("SELECT MAX(月份) FROM restaurants").fetchone()[0]

    stores = conn.execute("""
        SELECT 门店名称, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比, 堂食日均,
               机会点_消费人数不足, 机会点_人均频次不足, 机会点_超会不足, 机会点_好友拉新不足
        FROM restaurants
        WHERE 城市 = ? AND 月份 = ?
        ORDER BY 月消费人数 DESC
    """, (city_name, latest_month)).fetchall()

    city_metrics = conn.execute("""
        SELECT * FROM city_new_old WHERE 城市 = ? ORDER BY 月份 DESC LIMIT 1
    """, (city_name,)).fetchone()

    city_dinein = conn.execute("""
        SELECT * FROM city_dinein WHERE 城市 = ? ORDER BY 月份 DESC LIMIT 1
    """, (city_name,)).fetchone()

    trend = conn.execute("""
        SELECT 月份,
            ROUND(AVG(月消费人数), 0) as c,
            ROUND(AVG(人均月频次), 2) as f
        FROM restaurants
        WHERE 城市 = ?
        GROUP BY 月份 ORDER BY 月份
    """, (city_name,)).fetchall()

    conn.close()

    total = len(stores)
    opp_c = sum(1 for s in stores if s["机会点_消费人数不足"])
    opp_f = sum(1 for s in stores if s["机会点_人均频次不足"])
    opp_s = sum(1 for s in stores if s["机会点_超会不足"])
    opp_fr = sum(1 for s in stores if s["机会点_好友拉新不足"])

    return jsonify({
        "city": city_name,
        "store_count": total,
        "stores": [dict(s) for s in stores],
        "city_metrics": dict(city_metrics) if city_metrics else None,
        "city_dinein": dict(city_dinein) if city_dinein else None,
        "trend": {
            "months": [t["月份"] for t in trend],
            "consumers": [t["c"] for t in trend],
            "frequency": [t["f"] for t in trend],
        },
        "opportunity": {
            "consumer_pct": round(opp_c / total * 100, 1) if total else 0,
            "frequency_pct": round(opp_f / total * 100, 1) if total else 0,
            "super_pct": round(opp_s / total * 100, 1) if total else 0,
            "friend_pct": round(opp_fr / total * 100, 1) if total else 0,
        }
    })


# ═══════════════════════════════════════════
# User Management (Admin)
# ═══════════════════════════════════════════

@app.route("/api/users", methods=["GET"])
@login_required
def api_users_list():
    conn = get_db()
    users = conn.execute("SELECT id, username, display_name, role, created_at FROM users").fetchall()
    conn.close()
    return jsonify([dict(u) for u in users])


@app.route("/api/users", methods=["POST"])
@login_required
def api_users_create():
    data = request.get_json()
    username = data.get("username", "").strip()
    password = data.get("password", "")
    display_name = data.get("display_name", username)

    if not username or not password:
        return jsonify({"error": "用户名和密码不能为空"}), 400

    hashed = bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()

    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO users (username, password_hash, display_name) VALUES (?, ?, ?)",
            (username, hashed, display_name)
        )
        conn.commit()
    except sqlite3.IntegrityError:
        conn.close()
        return jsonify({"error": "用户名已存在"}), 409
    conn.close()

    return jsonify({"ok": True})


# ═══════════════════════════════════════════
# PPT Export
# ═══════════════════════════════════════════

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from urllib.parse import quote


@app.route("/api/store/<store_name>/ppt")
@login_required
def api_store_ppt(store_name):
    """Generate 2-page PPT report for a single store."""
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    conn = get_db()
    detail = conn.execute("""
        SELECT * FROM restaurants WHERE 门店名称 = ?
        ORDER BY 月份 DESC LIMIT 1
    """, (store_name,)).fetchone()

    if not detail:
        conn.close()
        return jsonify({"error": "未找到该餐厅"}), 404

    # Trend
    trend_rows = conn.execute("""
        SELECT 月份, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比,
               超会售卡, 机会点_消费人数不足, 机会点_人均频次不足,
               机会点_超会不足, 机会点_好友拉新不足, 堂食日均
        FROM restaurants WHERE 门店名称 = ? ORDER BY 月份
    """, (store_name,)).fetchall()

    # Peer
    city_type = detail["城市类型"]
    month = detail["月份"]
    peer = conn.execute("""
        SELECT
            AVG(月消费人数) as avg_consumers,
            AVG(人均月频次) as avg_frequency,
            AVG(超会消费占比) as avg_super_ratio,
            AVG(好友消费占比) as avg_friend_ratio,
            COUNT(*) as peer_count
        FROM restaurants
        WHERE 城市类型 = ? AND 月份 = ? AND 门店名称 != ?
    """, (city_type, month, store_name)).fetchone()

    # Ranking
    peers = conn.execute("""
        SELECT 门店名称, 月消费人数, 人均月频次, 超会消费占比, 好友消费占比
        FROM restaurants WHERE 城市类型 = ? AND 月份 = ?
    """, (city_type, month)).fetchall()

    def get_rank(field):
        sorted_peers = sorted(peers, key=lambda x: x[field] or 0, reverse=True)
        for i, p in enumerate(sorted_peers):
            if p["门店名称"] == store_name:
                return i + 1, len(sorted_peers)
        return None, len(sorted_peers)

    rank_c, total = get_rank("月消费人数")
    rank_f, _ = get_rank("人均月频次")
    rank_s, _ = get_rank("超会消费占比")
    rank_fr, _ = get_rank("好友消费占比")

    # Summary
    summary = generate_store_summary(store_name)
    conn.close()

    # Build PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    def add_text(slide, l, t, w, h, text, size=14, bold=False, color=RGBColor(0x33,0x33,0x33), align=None, font='Microsoft YaHei'):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        if isinstance(text, str): text = [text]
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            for r in p.runs:
                r.font.name = font
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color
            if align: p.alignment = align
        return tb

    def add_rect(slide, l, t, w, h, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(0.5)
        return shape

    def fmt_n(v, unit=''):
        if v is None: return '—'
        if unit=='万': return f'{v/10000:.1f}万'
        if v >= 100000000: return f'{v/100000000:.2f}亿'
        if v >= 10000: return f'{v/10000:.1f}万'
        return f'{round(v):,}'

    def fmt_p(v):
        if v is None: return '—'
        return f'{v*100:.1f}%'

    def fmt_f(v):
        if v is None: return '—'
        return f'{v:.2f}'

    blank = prs.slide_layouts[6]

    # ───────── Page 1 ─────────
    s1 = prs.slides.add_slide(blank)

    # Header
    add_rect(s1, 0, 0, SW, Inches(0.7), RGBColor(0,0x80,0x40))
    add_text(s1, Inches(0.4), Inches(0.15), Inches(12), Inches(0.5),
             f'ICU餐厅会员看板 · {detail["门店名称"]} · 单店分析',
             size=18, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s1, Inches(10.5), Inches(0.22), Inches(2.5), Inches(0.4),
             f'数据月份: {detail["月份"]}', size=11, color=RGBColor(0xE6,0xF5,0xEC))

    # Basic info row
    add_text(s1, Inches(0.4), Inches(0.95), Inches(12), Inches(0.4),
             f'📍 {detail["城市"] or ""} · {detail["战队"] or ""} · {detail["区域"] or ""} · {detail["城市类型"] or ""} · {detail["商圈类型"] or ""} · 开业:{detail["开业时间"] or ""}',
             size=12, color=RGBColor(0x77,0x77,0x77))

    # Section: 核心指标
    add_text(s1, Inches(0.4), Inches(1.55), Inches(12), Inches(0.4),
             '一、核心经营指标（最新月）', size=15, bold=True, color=RGBColor(0,0x80,0x40))

    # 5 metric cards
    metrics = [
        ('月消费人数', detail['月消费人数'], peer['avg_consumers'] if peer else None, fmt_n, '人'),
        ('人均月频次', detail['人均月频次'], peer['avg_frequency'] if peer else None, fmt_f, '次'),
        ('超会消费占比', detail['超会消费占比'], peer['avg_super_ratio'] if peer else None, fmt_p, ''),
        ('好友消费占比', detail['好友消费占比'], peer['avg_friend_ratio'] if peer else None, fmt_p, ''),
        ('超会售卡', detail['超会售卡'], None, fmt_n, '张'),
    ]
    card_w = Inches(2.4)
    card_h = Inches(1.4)
    start_x = Inches(0.4)
    gap = Inches(0.1)
    for i, (label, val, pv, fmt, unit) in enumerate(metrics):
        x = start_x + (card_w + gap) * i
        add_rect(s1, x, Inches(2.05), card_w, card_h, RGBColor(0xFA,0xFA,0xFA), line=RGBColor(0xE0,0xE0,0xE0))
        add_text(s1, x, Inches(2.1), card_w, Inches(0.3), label, size=11, color=RGBColor(0x88,0x88,0x88), align=2)
        add_text(s1, x, Inches(2.5), card_w, Inches(0.6), fmt(val) + (unit if val is not None else ''), size=22, bold=True, color=RGBColor(0x22,0x22,0x22), align=2)
        if pv is not None:
            diff = (val or 0) - pv
            pct_diff = (diff / pv * 100) if pv else 0
            color = RGBColor(0xE7,0x4C,0x3C) if diff > 0 else RGBColor(0x27,0xAE,0x60)
            sign = '+' if diff > 0 else ''
            add_text(s1, x, Inches(3.05), card_w, Inches(0.3),
                     f'同类均值 {fmt(pv)}  {sign}{pct_diff:.1f}%',
                     size=9, color=color, align=2)

    # Section: 同类排名
    add_text(s1, Inches(0.4), Inches(3.7), Inches(12), Inches(0.4),
             '二、同类门店排名', size=15, bold=True, color=RGBColor(0,0x80,0x40))

    rank_data = [
        ('月消费人数', rank_c, total),
        ('人均月频次', rank_f, total),
        ('超会消费占比', rank_s, total),
        ('好友消费占比', rank_fr, total),
    ]
    rank_w = Inches(2.95)
    rank_h = Inches(1.5)
    for i, (label, rank, tot) in enumerate(rank_data):
        x = Inches(0.4) + (rank_w + Inches(0.1)) * i
        add_rect(s1, x, Inches(4.15), rank_w, rank_h, RGBColor(0xFA,0xFA,0xFA), line=RGBColor(0xE0,0xE0,0xE0))
        add_text(s1, x, Inches(4.2), rank_w, Inches(0.3), label, size=11, color=RGBColor(0x88,0x88,0x88), align=2)
        pct = (rank / tot * 100) if rank and tot else None
        color = RGBColor(0x27,0xAE,0x60) if (pct and pct <= 50) else RGBColor(0xE7,0x4C,0x3C)
        add_text(s1, x, Inches(4.55), rank_w, Inches(0.7), str(rank or '—'), size=36, bold=True, color=color, align=2)
        add_text(s1, x, Inches(5.25), rank_w, Inches(0.4),
                 f'前{pct:.0f}% / 共{tot}家' if pct else '—',
                 size=11, color=color, align=2, bold=True)

    # Section: 经营诊断
    add_text(s1, Inches(0.4), Inches(5.85), Inches(12), Inches(0.4),
             '三、经营诊断', size=15, bold=True, color=RGBColor(0,0x80,0x40))
    if summary and summary.get('sentences'):
        for i, sent in enumerate(summary['sentences']):
            add_text(s1, Inches(0.4), Inches(6.3 + i*0.32), Inches(12.5), Inches(0.3),
                     f'{i+1}. {sent}', size=10.5, color=RGBColor(0x55,0x55,0x55))

    # Footer
    add_text(s1, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             f'ICU餐厅会员看板 · 数据截至 {detail["月份"]} · 第 1 / 3 页',
             size=9, color=RGBColor(0xAA,0xAA,0xAA), align=2)

    # ───────── Page 2 ─────────
    s2 = prs.slides.add_slide(blank)

    add_rect(s2, 0, 0, SW, Inches(0.7), RGBColor(0,0x80,0x40))
    add_text(s2, Inches(0.4), Inches(0.15), Inches(12), Inches(0.5),
             f'ICU餐厅会员看板 · {detail["门店名称"]} · 趋势追踪',
             size=18, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s2, Inches(10.5), Inches(0.22), Inches(2.5), Inches(0.4),
             f'共{len(trend_rows)}个月', size=11, color=RGBColor(0xE6,0xF5,0xEC))

    # Section: 月度趋势表
    add_text(s2, Inches(0.4), Inches(0.95), Inches(12), Inches(0.4),
             '一、10个月经营数据明细', size=15, bold=True, color=RGBColor(0,0x80,0x40))

    # Table
    table_data = [
        ['月份', '消费人数', '人均频次', '超会占比', '好友占比', '超会售卡', '堂食日均', '消费', '频次', '超会', '好友']
    ]
    for r in trend_rows:
        table_data.append([
            r['月份'],
            fmt_n(r['月消费人数']),
            fmt_f(r['人均月频次']),
            fmt_p(r['超会消费占比']),
            fmt_p(r['好友消费占比']),
            fmt_n(r['超会售卡']),
            fmt_n(r['堂食日均']),
            '⚠' if r['机会点_消费人数不足'] else '✓',
            '⚠' if r['机会点_人均频次不足'] else '✓',
            '⚠' if r['机会点_超会不足'] else '✓',
            '⚠' if r['机会点_好友拉新不足'] else '✓',
        ])

    rows, cols = len(table_data), len(table_data[0])
    table_shape = s2.shapes.add_table(rows, cols, Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.6))
    table = table_shape.table
    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = 2 if c_idx > 0 else 1
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(9)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    elif c_idx >= 7:  # opportunity columns
                        run.font.color.rgb = RGBColor(0xE7,0x4C,0x3C) if val == '⚠' else RGBColor(0x27,0xAE,0x60)
                    else:
                        run.font.color.rgb = RGBColor(0x33,0x33,0x33)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0,0x80,0x40)
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8,0xF8,0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    # Set column widths
    table.columns[0].width = Inches(0.9)
    for c in range(1, 7):
        table.columns[c].width = Inches(1.4)
    for c in range(7, 11):
        table.columns[c].width = Inches(0.6)

    # Section: 机会点变化
    add_text(s2, Inches(0.4), Inches(4.2), Inches(12), Inches(0.4),
             '二、机会点变化（红色=是，绿色=达标）', size=15, bold=True, color=RGBColor(0,0x80,0x40))

    # Draw opportunity grid
    grid_top = Inches(4.7)
    grid_left = Inches(0.4)
    cell_w = Inches(1.1)
    cell_h = Inches(0.5)
    label_w = Inches(1.5)
    opps = ['消费人数','消费频次','超会占比','好友拉新']
    opp_keys = [('机会点_消费人数不足',), ('机会点_人均频次不足',), ('机会点_超会不足',), ('机会点_好友拉新不足',)]

    # Header row
    add_rect(s2, grid_left, grid_top, label_w, cell_h, RGBColor(0xF0,0xF0,0xF0))
    add_text(s2, grid_left, grid_top, label_w, cell_h, '指标 / 月份', size=10, bold=True, align=2)
    for i, m in enumerate([r['月份'] for r in trend_rows]):
        x = grid_left + label_w + cell_w * i
        add_rect(s2, x, grid_top, cell_w, cell_h, RGBColor(0xF0,0xF0,0xF0))
        add_text(s2, x, grid_top, cell_w, cell_h, m, size=9, align=2, bold=True)

    # Data rows
    for r_idx, (opp_label, (key,)) in enumerate(zip(opps, opp_keys)):
        y = grid_top + cell_h * (r_idx + 1)
        add_rect(s2, grid_left, y, label_w, cell_h, RGBColor(0xF8,0xF8,0xF8))
        add_text(s2, grid_left, y, label_w, cell_h, opp_label, size=10, align=2)
        for c_idx, r in enumerate(trend_rows):
            x = grid_left + label_w + cell_w * c_idx
            val = r[key]
            color = RGBColor(0xE7,0x4C,0x3C) if val else RGBColor(0xE8,0xF5,0xE9)
            text_color = RGBColor(0xFF,0xFF,0xFF) if val else RGBColor(0x27,0xAE,0x60)
            add_rect(s2, x, y, cell_w, cell_h, color)
            add_text(s2, x, y, cell_w, cell_h, '⚠ 是' if val else '✓ 否', size=10, bold=True, color=text_color, align=2)

    add_text(s2, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             f'ICU餐厅会员看板 · 数据截至 {detail["月份"]} · 第 2 / 3 页',
             size=9, color=RGBColor(0xAA,0xAA,0xAA), align=2)

    # ───────── Page 3: 最新月机会点明细（Excel复刻表）─────────
    s3 = prs.slides.add_slide(blank)

    add_rect(s3, 0, 0, SW, Inches(0.7), RGBColor(0,0x80,0x40))
    add_text(s3, Inches(0.4), Inches(0.15), Inches(12), Inches(0.5),
             f'ICU餐厅会员看板 · {detail["门店名称"]} · 最新月机会点明细',
             size=18, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s3, Inches(10.5), Inches(0.22), Inches(2.5), Inches(0.4),
             f'数据月份: {detail["月份"]}', size=11, color=RGBColor(0xE6,0xF5,0xEC))

    add_text(s3, Inches(0.4), Inches(0.95), Inches(12), Inches(0.4),
             '最新月指标对照表（门店 vs 同类均值）', size=15, bold=True, color=RGBColor(0,0x80,0x40))

    snap_dims = [
        ('消费人数', detail['月消费人数'], detail.peer_consumers, detail['机会点_消费人数不足'], 'n'),
        ('人均月频次', detail['人均月频次'], detail.peer_frequency, detail['机会点_人均频次不足'], 'f'),
        ('新客复购率', detail['新客复购率'], None, None, 'p'),
        ('顾客复购率', detail['顾客复购率'], None, None, 'p'),
        ('超会消费占比', detail['超会消费占比'], detail.peer_super_ratio, detail['机会点_超会不足'], 'p'),
        ('超会售卡', detail['超会售卡'], None, None, 'n'),
        ('周四平均售卡', detail['周四平均售卡'], None, None, 'n'),
        ('好友消费占比', detail['好友消费占比'], detail.peer_friend_ratio, detail['机会点_好友拉新不足'], 'p'),
        ('好友新增', detail['好友新增'], None, None, 'n'),
        ('好友现存人数', detail['好友现存人数'], None, None, 'n'),
    ]

    def fmt_v(v, t):
        if v is None: return '—'
        if t=='p': return f'{v*100:.1f}%'
        if t=='f': return f'{v:.2f}'
        return f'{round(v):,}'

    snap_rows = [['指标', '本店', '同类均值', '机会点']]
    for label, val, peer, opp, t in snap_dims:
        snap_rows.append([label, fmt_v(val, t), fmt_v(peer, t),
                          '⚠ 是' if opp==1 else ('✓ 否' if opp==0 else '—')])

    rows, cols = len(snap_rows), len(snap_rows[0])
    snap_table = s3.shapes.add_table(rows, cols, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.5)).table
    snap_table.columns[0].width = Inches(3.0)
    snap_table.columns[1].width = Inches(3.0)
    snap_table.columns[2].width = Inches(3.0)
    snap_table.columns[3].width = Inches(3.5)

    for r_idx, row in enumerate(snap_rows):
        for c_idx, val in enumerate(row):
            cell = snap_table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = 2 if c_idx > 0 else 1
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(12)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    elif c_idx == 3:  # opportunity col
                        run.font.color.rgb = RGBColor(0xE7,0x4C,0x3C) if val == '⚠ 是' else RGBColor(0x27,0xAE,0x60)
                        run.font.bold = True
                    else:
                        run.font.color.rgb = RGBColor(0x33,0x33,0x33)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0,0x80,0x40)
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8,0xF8,0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
            cell.vertical_anchor = 3  # MIDDLE

    add_text(s3, Inches(0.4), Inches(6.1), Inches(12), Inches(0.4),
             '机会点判定：本店低于同类均值（按对应维度）即为「⚠ 是」，需关注。',
             size=11, color=RGBColor(0x55,0x55,0x55))

    add_text(s3, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             f'ICU餐厅会员看板 · 数据截至 {detail["月份"]} · 第 3 / 3 页',
             size=9, color=RGBColor(0xAA,0xAA,0xAA), align=2)

    # Save to buffer
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    from flask import send_file
    filename = f"{detail['门店名称']}_单店分析_{detail['月份']}.pptx"
    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


@app.route("/health")
def health():
    return "OK"

if __name__ == "__main__":
    print("=" * 50)
    print("ICU餐厅会员看板")
    print("默认账号: admin / admin123")
    print("访问地址: http://localhost:5000")
    print("=" * 50)
    app.run(host="0.0.0.0", port=5000, debug=True)
